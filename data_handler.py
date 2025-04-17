"""
Class that manages data for the RAG LLM. It prepares the data for the model, which should be populated in the data folder.

The data that will be accepeted in the following formats:
- PDF
- DOCX
- TXT
- PPTX
- NXML (specifically from statpearls)

Other formats will be ignored. Once cleaned the data will be extracted to a unified format (txt) the vectorized to be loaded into a vector DB.
"""

# Standard imports
import os
import re
import json
from pathlib import Path
from typing import Tuple, List

# Internal imports
from utils import embed_text
from const import PATH_TO_DATA, PATH_TO_CLEANED_DATA, PATH_TO_VECTORIZED_DATA

# External imports
import fitz
import docx2txt
import numpy as np
from tqdm import tqdm
import lxml.etree as et
from pptx import Presentation


class DataHandler:
    """
    Class that handles the data for the RAG LLM.

    Attributes:
    - data_path: str, path to the data folder
    - clean_data_path: str, path to the cleaned data folder
    - vectorized_data_path: str, path to the vectorized data folder
    - max_size_per_file: int, maximum size (number of key value pairs) of the dictionary before saving

    Methods:
    - load_data: loads the data from the data folder
    - clean_data: cleans the data
    """

    def __init__(
        self,
        data_path: Path = Path(PATH_TO_DATA),
        clean_data_path: Path = Path(PATH_TO_CLEANED_DATA),
        vectorized_data_path: Path = Path(PATH_TO_VECTORIZED_DATA),
        max_size_per_file: int = 1000,
    ) -> None:
        # Ensure the params are paths
        if not isinstance(data_path, Path):
            raise ValueError("Data path must be a Path object.")
        if not isinstance(clean_data_path, Path):
            raise ValueError("Clean data path must be a Path object.")
        if not isinstance(vectorized_data_path, Path):
            raise ValueError("Vectorized data path must be a Path object.")
        if not data_path.exists():
            raise FileNotFoundError(f"Data path {data_path} does not exist.")
        if not clean_data_path.exists():
            raise FileNotFoundError(
                f"Clean data path {clean_data_path} does not exist."
            )
        if not vectorized_data_path.exists():
            raise FileNotFoundError(
                f"Vectorized data path {vectorized_data_path} does not exist."
            )
        self.data_path = data_path
        self.clean_data_path = clean_data_path
        self.vectorized_data_path = vectorized_data_path
        self.data = []
        self.file_types = ["pdf", "docx", "txt", "pptx", "nxml"]
        self.max_size_per_file = (
            max_size_per_file  # Maximum size of the dictionary before saving
        )
        self.data_dict = {}  # Dictionary of titles and extracted content
        self.vectorized_data = {}  # Dictionary of titles and vectorized content

    def load_data(self) -> None:
        print("Loading data...")
        for root, dirs, files in os.walk(self.data_path):
            print(f"Reading {root}...")
            for file in tqdm(files):
                if file.split(".")[-1] in self.file_types:
                    self.data.append(os.path.join(root, file))
                # If .gitkeep file, ignore
                elif file == ".gitkeep":
                    continue
                # Else delete the file
                else:
                    print(f"File {file} is not in the accepted file types. Deleting...")
                    os.remove(os.path.join(root, file))

    def clean_data(self) -> None:
        print("Cleaning data...")
        for file in tqdm(self.data):
            if file.split(".")[-1] in ["pdf", "PDF"]:
                section_list = self.__clean_pdf(file)
                for title, text in section_list:
                    self.data_dict[title] = text
                    self.__write_to_file(self.clean_data_path, title, text)
                    continue
            elif file.split(".")[-1] in ["txt", "TXT"]:
                title, text = self.__clean_txt(file)
            elif file.split(".")[-1] in ["docx", "DOCX"]:
                title, text = self.__clean_docx(file)
            elif file.split(".")[-1] in ["pptx", "PPTX"]:
                title, text = self.__clean_pptx(file)
            elif file.split(".")[-1] == "nxml":
                title, text = self.__clean_nxml(file)
            else:
                print(
                    f"File {file} is not in the accepted file types. Should have been deleted... Skipping..."
                )
                continue

            if title is None or text is None:
                print(f"Error cleaning {file}. Skipping...")
                continue
            self.data_dict[title] = text

            self.__write_to_file(self.clean_data_path, title, text)

    def __write_to_file(self, path: Path, title: str, text: str) -> None:
        """
        Function that writes the cleaned data to a file.

        Parameters:
        - path: Path, path to the folder where the cleaned data will be saved
        - title: str, title of the data
        - text: str, cleaned text of the data
        """

        # Sanitize the title by replacing path separators and invalid filename characters
        invalid_chars = '<>:"/\\|?*,.\r\n'
        safe_title = "".join(c for c in title if c not in invalid_chars)

        # Remove BOM characters (if any)
        safe_title = safe_title.replace("\ufeff", "").replace("\ufffd", "")

        # Ensure safe title is less than 150 characters otherwise there is a risk of exceeding the limit
        safe_title = safe_title[:150]

        # Ensure that the title is not empty
        if not safe_title:
            safe_title = "Untitled"

        # Clean up any leading or trailing spaces and ensure the title is safe
        safe_title = safe_title.strip()

        # Make sure the path exists
        path.mkdir(parents=True, exist_ok=True)

        # Define the full path to save the cleaned data
        file_path = path / Path(safe_title + ".txt")

        # Save the cleaned data
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)

    def vectorize_data(self) -> None:
        """
        Function that vectorizes the data and saves them in multiple files when the dictionary size exceeds the limit.
        """
        if not self.data_dict:
            print(
                "Didn't clean data, assuming it's done already and saved. Loading data..."
            )
            # Get data from the cleaned data folder
            for root, dirs, files in os.walk(self.clean_data_path):
                print(f"Reading {root}...")
                for file in tqdm(files):
                    with open(os.path.join(root, file), "r", encoding="utf-8") as f:
                        text = f.read()
                    title = Path(file).stem
                    self.data_dict[title] = text

        if not self.data_dict:
            raise ValueError("No data to vectorize.")

        print("Vectorizing data...")

        current_count = 0  # To keep track of how many items are processed before saving
        file_counter = 1  # To keep track of file names

        # Vectorize the data and store them in batches
        for title, text in tqdm(self.data_dict.items()):
            if title == ".gitkeep":
                continue

            embeddings, chunks = embed_text(text)

            # Vectorize the text (assuming embed_text is defined elsewhere)
            self.vectorized_data[title] = {
                "embeddings": embeddings.tolist(),
                "texts": chunks,
            }

            current_count += 1  # Increment the count

            # Check if the current dictionary has reached the size limit
            if current_count >= self.max_size_per_file:
                # Save the vectorized data to a file
                self.__save_vectorized_data(file_counter)
                file_counter += 1  # Increment the file counter
                current_count = 0  # Reset the count
                self.vectorized_data = {}  # Clear the current dictionary to start fresh

        # Save any remaining data that was not saved in the last file
        if self.vectorized_data:
            self.__save_vectorized_data(file_counter)

    def __save_vectorized_data(self, file_counter: int) -> None:
        """
        Helper function to save vectorized data to a file.
        """
        file_name = f"vectorized_data_{file_counter}.json"
        with open(
            self.vectorized_data_path / Path(file_name), "w", encoding="utf-8"
        ) as f:
            json.dump(self.vectorized_data, f)
        # print(f"Saved vectorized data to {file_name}")

    def load_vectorized_data(self):
        """
        Generator that yields (key, vector) pairs from vectorized data files,
        instead of loading everything into memory at once.
        """
        for file in os.listdir(self.vectorized_data_path):
            if file.endswith(".json"):
                print(f"Loading {file}...")
                with open(
                    self.vectorized_data_path / Path(file), "r", encoding="utf-8"
                ) as f:
                    data = json.load(f)
                    for k, v in data.items():
                        yield k, dict(v)

    def __clean_pdf(self, file: str) -> List[Tuple[str, str]]:
        """
        Extracts text from a PDF file, splitting it by sections based on the Table of Contents (TOC),
        but only considers sections up to level 2 in the TOC.

        Parameters:
        - file: str, path to the PDF file.

        Returns:
        - List of tuples in the format ("title: section_title", text)
        """
        print(f"Cleaning {file}...")

        # Ensure it's a PDF
        if not file.lower().endswith(".pdf"):
            raise ValueError(f"File {file} is not a PDF file. Cannot clean as PDF.")

        pdf = fitz.open(file)
        title = Path(file).stem  # Extract filename without extension
        sections_list = []

        # Step 1: Extract TOC (Table of Contents)
        toc = pdf.get_toc()  # Returns a list: [(level, title, page), ...]

        if not toc:
            print("No TOC detected. Falling back to font-size based sectioning.")
            return self.__extract_pdf_by_font_size(file)  # Fallback method if no TOC

        # Step 2: Filter TOC to only include sections up to level 2
        toc = [entry for entry in toc if entry[0] <= 2]  # Limit to level 1 and 2

        # Step 3: Parse TOC into a dictionary {section_title: start_page}
        section_map = {
            entry[1]: entry[2] - 1 for entry in toc
        }  # TOC uses 1-based indexing

        # Step 4: Extract text per section
        section_titles = list(section_map.keys())  # Ordered list of sections
        for i, section_title in enumerate(section_titles):
            start_page = section_map[section_title]
            end_page = (
                section_map[section_titles[i + 1]]
                if i + 1 < len(section_titles)
                else len(pdf)
            )

            # Extract text from section pages
            section_text = []
            for page_num in range(start_page, end_page):
                section_text.append(pdf[page_num].get_text("text"))

            section_text = self.__clean_pdf_section_text(
                "\n".join(section_text).strip()
            )

            sections_list.append((f"{title}: {section_title}", section_text))

        pdf.close()
        return sections_list

    def __extract_pdf_by_font_size(self, file: str) -> list[Tuple[str, str]]:
        """
        Extracts text from a PDF file and splits it by detected chapter titles.

        Parameters:
        - file: str, path to the PDF file.

        Returns:
        - sections_list: list of tuples in the format ("title: section_title", text)
        """
        print(f"Cleaning {file} by font size since no TOC...")

        # Ensure it's a PDF
        if file.split(".")[-1].lower() != "pdf":
            raise ValueError(f"File {file} is not a PDF file. Cannot clean as PDF.")

        pdf = fitz.open(file)
        sections_list = []
        current_section = "Introduction"  # Default section title if no headers detected
        title = Path(file).stem  # Extracts filename without extension
        section_text = []

        for page in pdf:
            blocks = page.get_text("dict")["blocks"]
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        text = " ".join(
                            [
                                span["text"]
                                for span in line["spans"]
                                if span["text"].strip()
                            ]
                        )
                        if not text:
                            continue

                        # Detect possible section titles based on large/bold font sizes
                        for span in line["spans"]:
                            font_size = span["size"]

                            if font_size > 25:  # Tune thresholds as needed
                                # Store previous section
                                if section_text:
                                    sections_list.append(
                                        (
                                            f"{title}: {current_section}",
                                            self.__clean_pdf_section_text(
                                                "\n".join(section_text).strip()
                                            ),
                                        )
                                    )
                                    section_text = []

                                # Update section title
                                current_section = text.strip()
                                break  # Stop checking other spans if a title is found

                        section_text.append(text)

        # Append last section
        if section_text:
            sections_list.append(
                (f"{title}: {current_section}", "\n".join(section_text).strip())
            )

        pdf.close()
        return sections_list

    def __clean_pdf_section_text(self, section_text: str) -> str:
        """
        Cleans the text of a section extracted from a PDF.

        Parameters:
        - section_text: str, text extracted from a PDF section

        Returns:
        - cleaned_text: str, cleaned text
        """
        # Fix hyphenated words (e.g., "micro-\nscope" -> "microscope")
        section_text = re.sub(r"(\w+)-\n(\w+)", r"\1\2", section_text)

        # Remove line breaks that occur mid-sentence or in the middle of a paragraph
        section_text = re.sub(r"(?<!\n)\n(?!\n|\s*\w)", " ", section_text)

        # Ensure that sentence-ending line breaks are preserved
        section_text = re.sub(r"([.!?])\n+", r"\1\n", section_text)

        # Clean extra spaces, removing excessive whitespaces or spaces around line breaks
        section_text = re.sub(r"\s+", " ", section_text).strip()

        return section_text

    def __clean_txt(self, file: str) -> Tuple[str, str]:
        """
        Cleans a TXT file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["txt", "TXT"]:
            raise ValueError(f"File {file} is not a TXT file. Cannot clean as TXT.")
        # Open the TXT file
        with open(file, "r") as f:
            text = f.read()
        title = Path(file).stem
        return title, text

    def __clean_docx(self, file: str) -> Tuple[str, str]:
        """
        Cleans a DOCX file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["docx", "DOCX"]:
            raise ValueError(f"File {file} is not a DOCX file. Cannot clean as DOCX.")
        # Open the DOCX file
        text = docx2txt.process(file)
        title = file.split("/")[-1].split(".")[0]
        return title, text

    def __clean_pptx(self, file: str) -> Tuple[str, str]:
        """
        Cleans a PPTX file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["pptx", "PPTX"]:
            raise ValueError(f"File {file} is not a PPTX file. Cannot clean as PPTX.")
        # Open the PPTX file
        ppt = Presentation(file)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        title = Path(file).stem
        return title, text

    def __clean_nxml(self, file):
        """
        Cleans an nxml file. Specifically built to work with statpearls nxml files.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] != "nxml":
            raise ValueError(f"File {file} is not an nxml file. Cannot clean as nxml.")

        # Set up file extraction
        tree = et.parse(file)
        root = tree.getroot()

        # Extract the title using title-group tag
        try:
            title = root.find(".//title-group/title").text
        except AttributeError:
            print(f"Error extracting title from {file}. Skipping...")
            return None, None

        # print(f"Topic {title}...")

        # Extract tags with the sec-type as long as it does not have the value of "Continuing Education Activity"
        text = ""
        for sec in root.findall(".//sec"):
            sec_type = sec.get("sec-type")  # Access sec-type as an attribute of <sec>

            if sec_type is not None and sec_type != "Continuing Education Activity":
                for element in sec:  # Iterate over child elements in order
                    if element.tag == "title" and element.text:
                        text += element.text + "\n"
                    elif element.tag == "p" and element.text:
                        text += element.text + "\n"

        if text == "":
            print(f"Error extracting text from {file}. Skipping...")
            return None, None

        return f"StatPearls Chapter: {title}", text
